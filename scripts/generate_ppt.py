#!/usr/bin/env python
"""
Simple PPTX generator for Smart Study Partner.

Usage:
    python scripts\generate_ppt.py

This script reads `presentation_outline.md` in the project root and generates
`Smart_Study_Partner.pptx` in the project root. It requires the `python-pptx`
package: `pip install python-pptx`.
"""
from pathlib import Path
import sys


def read_outline(path: Path) -> list:
    text = path.read_text(encoding="utf-8")
    # Slides are separated by lines that contain only '---'
    sections = [s.strip() for s in text.split('\n---\n') if s.strip()]
    return sections


def make_pptx(sections, out_path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as e:
        print("Missing dependency: python-pptx. Install with: pip install python-pptx")
        raise

    prs = Presentation()

    # Title slide from the first section (if starts with 'Title:')
    first = sections[0] if sections else ''
    title = ''
    subtitle = ''
    for line in first.splitlines():
        if line.lower().startswith('title:'):
            title = line.split(':', 1)[1].strip()
        elif line.lower().startswith('subtitle:'):
            subtitle = line.split(':', 1)[1].strip()

    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = title
        if subtitle and len(slide.shapes) > 1:
            try:
                slide.placeholders[1].text = subtitle
            except Exception:
                pass

    # Remaining sections become slides
    for sec in sections[1:]:
        lines = [l.rstrip() for l in sec.splitlines() if l.strip()]
        if not lines:
            continue
        # Use first non-empty line as title
        s_title = lines[0].lstrip('# ').strip()
        body_lines = lines[1:]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s_title
        body = slide.shapes.placeholders[1].text_frame
        # Clear default paragraph
        body.clear()
        for bl in body_lines:
            # strip leading list markers
            text = bl.lstrip('-* ').strip()
            p = body.add_paragraph()
            p.text = text
            p.level = 0

    prs.save(str(out_path))


def main():
    root = Path(__file__).resolve().parents[1]
    outline = root / 'presentation_outline.md'
    if not outline.exists():
        print(f"Outline file not found: {outline}")
        sys.exit(2)

    sections = read_outline(outline)
    out_ppt = root / 'Smart_Study_Partner.pptx'
    try:
        make_pptx(sections, out_ppt)
    except Exception:
        sys.exit(1)

    print(f"Presentation generated: {out_ppt}")


if __name__ == '__main__':
    main()
